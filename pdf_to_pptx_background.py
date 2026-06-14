#!/usr/bin/env python
"""Convert PDF pages to a PPTX whose slides use the pages as true backgrounds."""

from __future__ import annotations

import argparse
import io
import sys
from collections.abc import Callable
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement


def set_slide_picture_background(slide, image_data: bytes) -> None:
    """Embed image_data and use it as the slide's picture background."""
    _, relationship_id = slide.part.get_or_add_image_part(io.BytesIO(image_data))
    common_slide_data = slide._element.cSld

    for child in list(common_slide_data):
        if child.tag == qn("p:bg"):
            common_slide_data.remove(child)

    background = OxmlElement("p:bg")
    background_properties = OxmlElement("p:bgPr")
    picture_fill = OxmlElement("a:blipFill")
    picture_fill.set("dpi", "0")
    picture_fill.set("rotWithShape", "1")

    image_reference = OxmlElement("a:blip")
    image_reference.set(qn("r:embed"), relationship_id)

    stretch = OxmlElement("a:stretch")
    stretch.append(OxmlElement("a:fillRect"))

    picture_fill.append(image_reference)
    picture_fill.append(stretch)
    background_properties.append(picture_fill)
    background_properties.append(OxmlElement("a:effectLst"))
    background.append(background_properties)

    # p:bg must appear before p:spTree inside p:cSld.
    common_slide_data.insert(0, background)


def convert_pdf(
    pdf_path: Path,
    output_path: Path,
    resolution: int,
    first_page: int,
    page_count: int | None,
    quiet: bool = False,
    progress_callback: Callable[[int, int, int], None] | None = None,
) -> int:
    with fitz.open(pdf_path) as document:
        total_pages = document.page_count
        start_index = first_page - 1

        if total_pages == 0:
            raise ValueError("PDF contains no pages.")
        if not 0 <= start_index < total_pages:
            raise ValueError(
                f"--from-page must be between 1 and {total_pages}, got {first_page}."
            )
        if page_count is not None and page_count < 1:
            raise ValueError("--count must be at least 1.")

        end_index = total_pages
        if page_count is not None:
            end_index = min(start_index + page_count, total_pages)

        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]

        first_pdf_page = document.load_page(start_index)
        aspect_ratio = first_pdf_page.rect.width / first_pdf_page.rect.height
        presentation.slide_width = int(presentation.slide_height * aspect_ratio)

        matrix = fitz.Matrix(resolution / 72, resolution / 72)
        selected_pages = range(start_index, end_index)
        selected_count = len(selected_pages)

        for item_number, page_number in enumerate(selected_pages, start=1):
            page = document.load_page(page_number)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            png_data = pixmap.tobytes(output="png")

            slide = presentation.slides.add_slide(blank_layout)
            set_slide_picture_background(slide, png_data)

            if progress_callback is not None:
                progress_callback(item_number, selected_count, page_number + 1)

            if not quiet:
                print(
                    f"[{item_number}/{selected_count}] "
                    f"PDF page {page_number + 1} converted"
                )

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(output_path)
        return selected_count


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Convert PDF pages into a PPTX with each rendered page stored as "
            "the slide's true picture background."
        )
    )
    parser.add_argument("pdf", type=Path, help="Input PDF file")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        help="Output PPTX; defaults to <PDF name>-background.pptx",
    )
    parser.add_argument(
        "-r",
        "--resolution",
        type=int,
        default=300,
        help="Rendering resolution in DPI (default: 300)",
    )
    parser.add_argument(
        "--from-page",
        type=int,
        default=1,
        help="First PDF page to include, using 1-based numbering (default: 1)",
    )
    parser.add_argument(
        "--count",
        type=int,
        default=None,
        help="Number of pages to include (default: all remaining pages)",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Allow replacing an existing output PPTX",
    )
    parser.add_argument(
        "-q",
        "--quiet",
        action="store_true",
        help="Suppress per-page progress messages",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    pdf_path = args.pdf.expanduser().resolve()
    output_path = (
        args.output.expanduser().resolve()
        if args.output
        else pdf_path.with_name(f"{pdf_path.stem}-background.pptx")
    )

    if not pdf_path.is_file():
        print(f"Input PDF not found: {pdf_path}", file=sys.stderr)
        return 2
    if args.resolution < 72:
        print("--resolution must be at least 72 DPI.", file=sys.stderr)
        return 2
    if output_path.exists() and not args.overwrite:
        print(
            f"Output already exists: {output_path}\n"
            "Use --overwrite or choose another path with --output.",
            file=sys.stderr,
        )
        return 2

    try:
        slide_count = convert_pdf(
            pdf_path=pdf_path,
            output_path=output_path,
            resolution=args.resolution,
            first_page=args.from_page,
            page_count=args.count,
            quiet=args.quiet,
        )
    except Exception as exc:
        print(f"Conversion failed: {exc}", file=sys.stderr)
        return 1

    print(f"Created {output_path} ({slide_count} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
